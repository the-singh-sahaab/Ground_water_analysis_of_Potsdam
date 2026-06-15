from pptx import Presentation

# Load reference PPT
ref_prs = Presentation('presentations/Hydro_Research_Pressentation (004).pptx')
print(f'Reference PPT has {len(ref_prs.slides)} slides')
print(f'Slide width: {ref_prs.slide_width}, height: {ref_prs.slide_height}')

# Check first few slides
for i, slide in enumerate(list(ref_prs.slides)[:3]):
    print(f'\nSlide {i+1}:')
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text = shape.text.strip()
            if text:
                print(f'  - Text: {text[:50]}')
        if shape.shape_type == 13:  # Picture
            print(f'  - Has picture')
